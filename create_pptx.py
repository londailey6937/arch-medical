#!/usr/bin/env python3
"""Generate Control Tower Sales Pitch PowerPoint Presentation."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Paths ────────────────────────────────────────────────────────────────
SCREENSHOTS = "/Users/londailey/510kBridge_Webpage/public/screenshots"
OUTPUT = "/Users/londailey/510kBridge_Webpage/Control_Tower_Sales_Deck.pptx"

# ── Brand colours ────────────────────────────────────────────────────────
NAVY   = RGBColor(0x0B, 0x23, 0x3A)
TEAL   = RGBColor(0x00, 0x96, 0x88)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY  = RGBColor(0xF5, 0xF5, 0xF5)
DGRAY  = RGBColor(0x33, 0x33, 0x33)
ACCENT = RGBColor(0x00, 0x7A, 0xCC)
RED    = RGBColor(0xE5, 0x3E, 0x3E)
GREEN  = RGBColor(0x2E, 0x7D, 0x32)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height


# ── Helpers ──────────────────────────────────────────────────────────────
def add_bg(slide, color=NAVY):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_bg(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_bullet_frame(slide, left, top, width, height, bullets, font_size=14,
                     color=DGRAY, title=None, title_size=18, title_color=NAVY):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    if title:
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(title_size)
        p.font.color.rgb = title_color
        p.font.bold = True
        p.font.name = "Calibri"
        p.space_after = Pt(8)
        start = 1
    else:
        start = 0

    for i, bullet in enumerate(bullets):
        if i == 0 and not title:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(4)
        p.level = 0
    return tf


def screenshot_path(name):
    path = os.path.join(SCREENSHOTS, f"{name}.png")
    return path if os.path.exists(path) else None


def add_screenshot(slide, name, left, top, width, height=None):
    path = screenshot_path(name)
    if not path:
        return None
    if height:
        pic = slide.shapes.add_picture(path, left, top, width, height)
    else:
        pic = slide.shapes.add_picture(path, left, top, width)
    return pic


def tab_slide(slide, tab_number, tab_name, group_name):
    """Common header band for every tab slide."""
    add_bg(slide, WHITE)
    # Top navy bar
    add_shape_bg(slide, Inches(0), Inches(0), W, Inches(1.1), NAVY)
    add_text_box(slide, Inches(0.6), Inches(0.15), Inches(8), Inches(0.5),
                 f"{group_name}",
                 font_size=13, color=TEAL, bold=True)
    add_text_box(slide, Inches(0.6), Inches(0.45), Inches(10), Inches(0.55),
                 f"Tab {tab_number}: {tab_name}",
                 font_size=26, color=WHITE, bold=True)
    # Branding
    add_text_box(slide, Inches(10.5), Inches(0.3), Inches(2.5), Inches(0.5),
                 "510(k) Bridge", font_size=16, color=TEAL, bold=True,
                 alignment=PP_ALIGN.RIGHT)
    add_text_box(slide, Inches(10.5), Inches(0.6), Inches(2.5), Inches(0.4),
                 "Control Tower", font_size=12, color=WHITE,
                 alignment=PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, NAVY)

add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1),
             "510(k) Bridge Control Tower",
             font_size=44, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(2.6), Inches(11), Inches(0.8),
             "The First Purpose-Built FDA 510(k) Project Management Platform",
             font_size=22, color=TEAL, bold=False, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(2), Inches(3.8), Inches(9), Inches(1.5),
             "16 integrated tabs  •  Trilingual (EN / 中文 / 한국어)  •  Role-based access\n"
             "21 CFR Part 11 audit trail  •  ISO 14971 risk management  •  FDA Q-Sub automation\n"
             "Real-time cloud sync  •  7 pre-configured device templates  •  3 subscription tiers",
             font_size=16, color=WHITE, alignment=PP_ALIGN.CENTER)

# Divider line
add_shape_bg(slide, Inches(4), Inches(5.6), Inches(5.3), Inches(0.03), TEAL)

add_text_box(slide, Inches(1), Inches(6.0), Inches(11), Inches(0.5),
             "510(k) Bridge  •  Gresham, Oregon  •  www.510kbridge.com",
             font_size=14, color=WHITE, alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 2 — MARKET PROBLEM
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), W, Inches(1.1), NAVY)
add_text_box(slide, Inches(0.6), Inches(0.3), Inches(10), Inches(0.6),
             "The Problem: No FDA-Specific Project Management Platform Exists",
             font_size=28, color=WHITE, bold=True)

add_bullet_frame(slide, Inches(0.8), Inches(1.4), Inches(5.8), Inches(5.5), [
    "▸ Medical device companies use generic PM tools (Jira, Asana, Monday.com) that have zero FDA awareness",
    "▸ Regulatory consultants rely on spreadsheets, email threads, and tribal knowledge",
    "▸ No existing platform connects technical milestones to FDA regulatory gates",
    "▸ 510(k) submissions average 6-9 months of rework due to disorganized documentation (DHF/DMR gaps)",
    "▸ Companies send 2-3 FDA Q-Sub requests before getting accepted — each cycle costs $10K-$25K in delays",
    "▸ ISO 14971 risk management is tracked in standalone documents, disconnected from project decisions",
    "▸ International teams (US-China, US-Korea corridors) lack a unified multilingual workspace",
], font_size=15, color=DGRAY, title="Why the Industry Needs Control Tower",
   title_size=22, title_color=NAVY)

add_bullet_frame(slide, Inches(7.0), Inches(1.4), Inches(5.8), Inches(5.5), [
    "▸ Greenlight Guru — QMS-focused, not project management; no milestone/gate tracking; $$$",
    "▸ MasterControl — Enterprise QMS; complex, expensive ($50K+/yr); no 510(k) workflow",
    "▸ Jama Connect — Requirements management only; no budget, risk matrix, or FDA comms",
    "▸ Arena Solutions (PTC) — PLM/BOM-focused; no regulatory submission tools",
    "▸ Qualio — Document-centric QMS; no dual-track milestones, no gate system",
    "▸ Monday.com / Asana — Generic PM; zero FDA vocabulary or compliance features",
    "",
    "None of these platforms provide an integrated 510(k) project lifecycle from "
    "Pre-Sub through clearance with built-in FDA automation, trilingual support, "
    "and role-based governance.",
], font_size=15, color=DGRAY, title="Competitive Landscape Gaps",
   title_size=22, title_color=RED)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 3 — PLATFORM OVERVIEW / ARCHITECTURE
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), W, Inches(1.1), NAVY)
add_text_box(slide, Inches(0.6), Inches(0.3), Inches(10), Inches(0.6),
             "Platform Architecture: 16 Integrated Tabs in 5 Groups",
             font_size=28, color=WHITE, bold=True)

groups = [
    ("PROGRAM MGMT", TEAL, ["1. Dual-Track Milestones", "2. Gate System",
                             "3. Timeline", "4. Actions Panel"]),
    ("REGULATORY", RGBColor(0xE6, 0x5C, 0x00), ["5. Regulatory Tracker",
                  "6. Risk Dashboard", "7. FDA Comms Center"]),
    ("DOCUMENTS", ACCENT, ["8. Document Control", "9. Audit Trail"]),
    ("FINANCE", RGBColor(0x7B, 0x1F, 0xA2), ["10. Budget", "11. Cash / Runway",
               "12. US Investment", "13. Cap Table"]),
    ("OPERATIONS", GREEN, ["14. Resources", "15. Suppliers",
                           "16. Message Board"]),
]

x_start = Inches(0.5)
y_start = Inches(1.5)
col_w = Inches(2.4)
gap = Inches(0.15)

for i, (name, color, tabs) in enumerate(groups):
    left = x_start + i * (col_w + gap)
    # Group header
    add_shape_bg(slide, left, y_start, col_w, Inches(0.5), color)
    add_text_box(slide, left, y_start + Inches(0.05), col_w, Inches(0.4),
                 name, font_size=13, color=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)
    # Tab list
    add_shape_bg(slide, left, y_start + Inches(0.5), col_w, Inches(3.0), LGRAY)
    add_bullet_frame(slide, left + Inches(0.15), y_start + Inches(0.6),
                     col_w - Inches(0.3), Inches(2.8), tabs,
                     font_size=13, color=DGRAY)

# Cross-cutting features
add_shape_bg(slide, Inches(0.5), Inches(5.3), Inches(12.3), Inches(0.05), TEAL)
add_bullet_frame(slide, Inches(0.5), Inches(5.5), Inches(12.3), Inches(1.5), [
    "Trilingual (EN/中文/한국어)  •  Role-Based Access (PMP / Tech / Business / Accounting)  •  "
    "3 Subscription Tiers ($500–$2,000/mo)  •  7 Device Templates  •  Supabase Cloud + Offline  •  "
    "21 CFR Part 11 Audit Trail  •  Setup Wizard  •  Real-Time Collaboration"
], font_size=14, color=DGRAY, title="Cross-Cutting Platform Capabilities",
   title_size=18, title_color=NAVY)


# ═══════════════════════════════════════════════════════════════════════
# TAB SLIDES — one per tab with screenshot + sales pitch + competitor gap
# ═══════════════════════════════════════════════════════════════════════

tab_data = [
    # (tab_num, tab_name, group, screenshot_key, value_props, competitor_gap)
    (1, "Dual-Track Milestones", "PROGRAM MANAGEMENT", "dual-track", [
        "▸ Parallel technical & regulatory milestone tracking — the ONLY platform that visualizes both tracks side-by-side",
        "▸ Technical track: Design freeze → Prototype → Bench verification → V&V → Design transfer",
        "▸ Regulatory track: Pre-Sub Q-Meeting → 510(k) prep → Submission → Expected clearance",
        "▸ Auto-populated timelines based on project duration (18-23 months typical)",
        "▸ Status cycling (Not Started → In Progress → Complete / Blocked) with full audit logging",
        "▸ Owner assignment per milestone (Tech, Regulatory, or Business) ensures accountability",
    ],
     "Jira/Asana have task boards but no concept of synchronized technical + regulatory tracks. "
     "Greenlight Guru tracks QMS documents but not engineering milestones. "
     "No competitor links a prototype testing schedule to an FDA submission timeline."
    ),

    (2, "Gate System", "PROGRAM MANAGEMENT", "gate-system", [
        "▸ Phase-gate governance with digital review — checklist-based criteria per gate",
        "▸ 2-6 gates auto-generated based on project length",
        "▸ Gate decisions: Proceed / No-Go / Conditional — PMP records with timestamp attribution",
        "▸ Stakeholder inputs from Tech and Business teams feed into each gate decision",
        "▸ Gate notes capture discussion points, conditions, and commitments",
        "▸ Gate status auto-updates (Not Started → Pending Review → Approved → Blocked) as criteria complete",
    ],
     "Stage-Gate® is a well-known methodology, but no SaaS tool implements it for medical devices. "
     "Generic PM tools require manual spreadsheet gates. "
     "Control Tower's gate system directly references 510(k) deliverables — "
     "e.g., 'DHF complete?' as an actual gate criterion linked to the DHF tracker."
    ),

    (3, "Timeline", "PROGRAM MANAGEMENT", "timeline", [
        "▸ Month-by-month consolidated project view showing technical + business activities",
        "▸ Impact indicators: Neutral (|), Warning (⚠), Critical (🔴) for schedule risk",
        "▸ Auto-generated from device template or wizard inputs — instant project skeleton",
        "▸ Editable events with team-modifiable text and impact levels",
        "▸ Events linked to Dual-Track milestones for full traceability",
        "▸ Gives executives a single-page view of the entire 510(k) program lifecycle",
    ],
     "Gantt charts in MS Project or Monday.com are generic bars — "
     "they don't know an FDA milestone from a marketing campaign. "
     "Control Tower's timeline is purpose-built for 510(k) programs "
     "with built-in FDA terminology and impact-flagging."
    ),

    (4, "Actions Panel", "PROGRAM MANAGEMENT", "actions", [
        "▸ Full task board: title, assignee, priority (H/M/L), status (Todo/In Progress/Done/Blocked), due date",
        "▸ Gate-linked actions — each task can reference a specific gate (G1, G2, etc.)",
        "▸ DHF Document Tracker — 21 CFR 820.30 Design History File document checklist",
        "▸ DMR Document Tracker — 21 CFR 820.181 Device Master Record (12 document types)",
        "▸ CAPA Log — Corrective & Preventive Actions with risk linkage from Risk Dashboard",
        "▸ Four integrated sub-sections in one tab: Tasks + DHF + DMR + CAPA",
    ],
     "Greenlight Guru has DHF tracking but it's siloed from project management tasks. "
     "Jira has task boards but no DHF/DMR/CAPA awareness. "
     "Control Tower uniquely combines task management with FDA documentation checklists "
     "and links CAPAs directly to risks identified in the Risk Dashboard."
    ),

    (5, "Regulatory Tracker", "REGULATORY", "regulatory-tracker", [
        "▸ Standards compliance dashboard auto-populated from device template",
        "▸ Tracks IEC 60601-1, IEC 60601-1-2, ISO 10993, ISO 14971, IEC 62304, 21 CFR 820, ISO 13485",
        "▸ Clause-level tracking — drill down to individual sub-clauses (e.g., ISO 10993-5, -10)",
        "▸ 0-100% progress bar per standard with status cycling",
        "▸ Evidence linking — attach document references to each clause",
        "▸ Instantly shows which standards are behind schedule, focusing team effort",
    ],
     "Qualio and MasterControl track document compliance but not standard-by-standard progress. "
     "No competitor auto-populates applicable standards based on device type. "
     "Teams currently maintain compliance matrices in Excel — "
     "Control Tower eliminates that spreadsheet and makes it live."
    ),

    (6, "Risk Dashboard", "REGULATORY", "risk-dashboard", [
        "▸ Full ISO 14971 risk management visualization with severity × probability matrix",
        "▸ Color-coded risk levels: Green (Acceptable) → Yellow (ALARP) → Red (Unacceptable)",
        "▸ Per-risk fields: severity, probability, controls, residual risk, mitigation status, module, standard",
        "▸ Red-risk alerts automatically trigger warnings in the FDA Comms tab",
        "▸ Template-pre-populated risks (5-8 per device category — e.g., respiratory gets EMG-specific risks)",
        "▸ Click risk ID to open detailed editor — all changes audit-logged",
    ],
     "Most companies track risks in Word documents or standalone Excel matrices. "
     "Greenlight Guru has risk management but decoupled from project timelines. "
     "Control Tower's Risk Dashboard connects risks to CAPAs, gate criteria, "
     "and FDA Pre-Sub questions — a closed-loop risk system."
    ),

    (7, "FDA Communications Center", "REGULATORY", "fda-comms", [
        "▸ PMP-exclusive tab — the nerve center for FDA interactions",
        "▸ Q-Sub Cover Letter Generator — 5 types (Pre-Sub Meeting, Written Questions, Controlled Correspondence, etc.)",
        "▸ Export Question Package — pulls Pre-Sub question threads from Message Board into formatted FDA package",
        "▸ Refuse-to-Accept (RTA) Checklist — 17-item self-validation against FDA acceptance criteria",
        "▸ MDUFA Review Timeline — Day 1 through Day 90 milestone tracking for standard 510(k) review",
        "▸ Substantial Equivalence Decision Flowchart — FDA's SE determination logic visualized",
    ],
     "NO COMPETITOR HAS THIS. No existing platform generates FDA Q-Sub cover letters, "
     "exports Pre-Sub question packages, or provides an RTA pre-check. "
     "This is Control Tower's single strongest differentiator — "
     "it automates the most time-consuming and error-prone regulatory communication workflows."
    ),

    (8, "Document Control", "DOCUMENTS", "document-control", [
        "▸ ISO 13485-aligned document lifecycle: Draft → In Review → Approved → Effective → Obsolete",
        "▸ Auto-generated Document Control Numbers (DCN-REG-001, etc.)",
        "▸ Categories: Regulatory, Technical, Business, Legal, Finance, Templates",
        "▸ Browser localStorage for IP protection + optional server sync for approved docs",
        "▸ Source reference linking (GitHub commit, SVN revision, etc.)",
        "▸ Milestone linkage — e.g., R8 milestone triggers DHF document readiness check",
    ],
     "MasterControl and Qualio are full QMS platforms costing $50K+/yr. "
     "Control Tower provides 80% of critical document control features "
     "at a fraction of the cost, tightly integrated with project milestones "
     "rather than existing as a standalone system."
    ),

    (9, "Audit Trail", "DOCUMENTS", "audit-trail", [
        "▸ 21 CFR Part 11-compliant immutable activity log",
        "▸ Every change captured: timestamp, user role, action type, target, old value → new value",
        "▸ 100+ action types (milestone-status, risk-field, gate-decision, doc-status, budget-entry, etc.)",
        "▸ Keyword search and action-type filtering",
        "▸ Export complete history as CSV/JSON for regulatory submissions",
        "▸ Auto-synced to Supabase backend; queues locally if offline — zero data loss",
    ],
     "This is table-stakes for regulated industries but missing from generic PM tools. "
     "Jira's audit log is IT-focused, not regulatory. "
     "Control Tower's audit trail is designed specifically for FDA investigators and notified body auditors, "
     "capturing the exact fields and formatting they expect."
    ),

    (10, "Budget", "FINANCE", "budget", [
        "▸ Budget vs. actual tracking by category — pre-populated from device template",
        "▸ Categories: Prototype & Materials, Testing, Regulatory & Legal, Personnel, etc.",
        "▸ Per-category fields: planned amount, actual amount, variance (auto-calculated), notes",
        "▸ Dual currency support: USD/CNY auto-toggle (7.25× conversion rate)",
        "▸ Add/edit/delete categories as project evolves",
        "▸ All budget changes captured in Audit Trail for financial accountability",
    ],
     "Generic PM tools either lack budgeting entirely (Jira) or provide basic spend tracking "
     "without medical-device-specific cost categories. "
     "Control Tower pre-populates budget categories that match 510(k) development phases — "
     "testing lab fees, regulatory counsel, biocompatibility studies, etc."
    ),

    (11, "Cash / Runway", "FINANCE", "cash-runway", [
        "▸ Real-time financial health dashboard: Cash on Hand, Monthly Burn Rate, Runway (months)",
        "▸ Runway color coding: Green (>12 mo), Yellow (6-12 mo), Red (<6 mo) — instant visibility",
        "▸ Funding round tracking: Pipeline → Committed → Received with amounts and dates",
        "▸ Month-by-month burn history: actual vs. planned expenditure",
        "▸ Visual burn history charts for board presentations",
        "▸ Critical for startups — answers 'do we have enough cash to reach FDA clearance?'",
    ],
     "Financial planning tools (QuickBooks, Xero) don't tie cash runway to regulatory milestones. "
     "Startup dashboards (Carta, Visible) track burn but know nothing about FDA timelines. "
     "Control Tower uniquely answers the question every med-device founder asks: "
     "'Will my money last until 510(k) clearance?'"
    ),

    (12, "US Investment", "FINANCE", "us-investment", [
        "▸ Investor relations pipeline: Prospect → Contacted → In Discussions → Term Sheet → Committed",
        "▸ Investor profiles: type (VC, Angel, Strategic, PE, Government), stage fit, committed amounts",
        "▸ IR activity logging: meetings, presentations, due diligence sessions, term negotiations",
        "▸ Activity metrics and conversion rates for fundraising effectiveness",
        "▸ Purpose-built for international med-device companies entering the US market",
        "▸ Consolidates investor tracking that otherwise lives in scattered spreadsheets and CRMs",
    ],
     "CRM tools (HubSpot, Salesforce) track investor contacts but have no medical-device context. "
     "Carta handles cap tables but not investor pipeline. "
     "Control Tower provides a specialized investor pipeline that sits alongside the 510(k) program — "
     "so founders can show investors exactly where they are in the regulatory process."
    ),

    (13, "Cap Table", "FINANCE", None, [
        "▸ Full equity ownership registry: Common, Preferred A/B/C, Options, Warrants",
        "▸ Auto-calculated ownership percentages and board seat tracking",
        "▸ Equity events: funding rounds, stock splits, option grants, preferred conversions, warrant exercises",
        "▸ Vesting schedule administration: 4-year standard with 1-year cliff (customizable)",
        "▸ Current vested/unvested tracking with next vest date calculation",
        "▸ Three integrated sections: Shareholders + Equity Events + Vesting Schedules",
    ],
     "Carta is the gold standard for cap tables but costs $5K-$50K/yr and is standalone. "
     "Control Tower's cap table is integrated with the project — "
     "when a funding round closes (US Investment tab), it can flow into the cap table. "
     "For early-stage med-device startups, this eliminates yet another standalone tool."
    ),

    (14, "Resources", "OPERATIONS", "resources", [
        "▸ Team allocation and utilization monitoring across workstreams",
        "▸ Workstream allocation: % across Technical, Regulatory, Business, Finance (+ custom)",
        "▸ Utilization gauge: Green (<85%), Yellow (85-100%), Red (>100%) — prevents burnout",
        "▸ Team member profiles: name, role, email, bio",
        "▸ Inline editing — click percentage to adjust allocations in real time",
        "▸ Audit-logged allocation changes for workforce planning transparency",
    ],
     "Resource management in Jira or Monday.com is task-hour based, not workstream-% based. "
     "Med-device teams often have dual-hatted engineers (tech + regulatory). "
     "Control Tower's allocation model matches how small device teams actually work — "
     "showing at a glance if your V&V engineer is also overloaded on regulatory documentation."
    ),

    (15, "Suppliers", "OPERATIONS", "suppliers", [
        "▸ Supplier qualification and PO tracking — aligned with 21 CFR 820 supplier controls",
        "▸ Status lifecycle: Under Review → Qualified → Active → On Hold → Rejected",
        "▸ Fields: component/part, lead time (days), PO status, contract mfg milestones",
        "▸ Qualification criteria: ISO 13485/9001 certs, test reports, business continuity plans",
        "▸ Critical for 510(k) — FDA expects documented supplier qualification in your QMS",
        "▸ Linked to project milestones — knows when a component's lead time threatens a gate review",
    ],
     "Arena Solutions (PTC) is a full PLM with BOM management but costs $1,000+/seat/yr. "
     "Control Tower provides the supplier qualification tracking FDA expects "
     "without the complexity of a full PLM — perfect for startups "
     "building their first device with 5-15 suppliers."
    ),

    (16, "Message Board", "COMMUNICATIONS", "message-board", [
        "▸ Cross-functional threaded discussions with decision & action tracking",
        "▸ Workstreams: Technical, Regulatory, Business, Finance, Pre-Sub Questions, Other",
        "▸ Intent tags: Discuss, Decide, Inform, Escalate — every conversation has purpose",
        "▸ [DECISION] and [ACTION] message prefixes for traceability",
        "▸ Pre-Sub Questions workstream: collaboratively author FDA questions → export to FDA Comms",
        "▸ Views: All Threads, My Items, Decisions, Executive — filter by lifecycle (Open/Resolved)",
    ],
     "Slack, Teams, and email are where medical device decisions die — "
     "critical regulatory decisions buried in chat streams with no traceability. "
     "Control Tower's Message Board captures every decision with role attribution, "
     "timestamps, and the ability to spawn action items — "
     "then exports Pre-Sub questions directly to the FDA Communications Center."
    ),
]

for tab_num, tab_name, group, screenshot_key, value_props, competitor_gap in tab_data:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tab_slide(slide, tab_num, tab_name, group)

    # Left column: screenshot
    ss = None
    if screenshot_key:
        ss = add_screenshot(slide, screenshot_key,
                            Inches(0.4), Inches(1.3), Inches(6.2))

    # If screenshot present, place text on right; otherwise full width
    if ss:
        text_left = Inches(6.9)
        text_width = Inches(6.0)
    else:
        text_left = Inches(0.8)
        text_width = Inches(11.5)

    # Value props
    add_bullet_frame(slide, text_left, Inches(1.3), text_width, Inches(3.5),
                     value_props, font_size=12, color=DGRAY,
                     title="Value Proposition", title_size=16, title_color=TEAL)

    # Competitor gap — bottom band
    add_shape_bg(slide, Inches(0.3), Inches(5.6), Inches(12.7), Inches(1.6), LGRAY)
    add_text_box(slide, Inches(0.5), Inches(5.65), Inches(3), Inches(0.35),
                 "Competitor Gap Analysis", font_size=13, color=RED, bold=True)
    add_text_box(slide, Inches(0.5), Inches(5.95), Inches(12.3), Inches(1.2),
                 competitor_gap, font_size=12, color=DGRAY)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 20 — DEVICE TEMPLATES & WIZARD
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), W, Inches(1.1), NAVY)
add_text_box(slide, Inches(0.6), Inches(0.3), Inches(10), Inches(0.6),
             "Setup Wizard & 7 Pre-Configured Device Templates",
             font_size=28, color=WHITE, bold=True)

templates = [
    ("🫁 Respiratory", "Ventilators, CPAP, nebulizers\nProduct codes: BZD, CBK, BTT, FRA"),
    ("❤️ Cardiovascular", "ECG monitors, BP devices, stents\nProduct codes: DRX, DXH, MHX"),
    ("🦴 Orthopedic", "Joint implants, instruments,\nfixation devices"),
    ("🧪 IVD", "In-vitro diagnostics, assays,\nanalyzers"),
    ("📷 Imaging", "X-ray, ultrasound,\nMRI accessories"),
    ("♿ Rehabilitation", "Therapy devices,\nmobility aids"),
    ("💻 SaMD", "Software as a Medical Device\n(IEC 62304 lifecycle)"),
]

for i, (name, desc) in enumerate(templates):
    col = i % 4
    row = i // 4
    left = Inches(0.5) + col * Inches(3.15)
    top = Inches(1.5) + row * Inches(2.5)
    add_shape_bg(slide, left, top, Inches(2.9), Inches(2.2), LGRAY)
    add_text_box(slide, left + Inches(0.15), top + Inches(0.1),
                 Inches(2.6), Inches(0.4), name,
                 font_size=16, color=NAVY, bold=True)
    add_text_box(slide, left + Inches(0.15), top + Inches(0.55),
                 Inches(2.6), Inches(0.7), desc,
                 font_size=11, color=DGRAY)
    add_text_box(slide, left + Inches(0.15), top + Inches(1.3),
                 Inches(2.6), Inches(0.8),
                 "Includes:\n• Pre-seeded risks & standards\n• Budget categories\n• Timeline & milestones",
                 font_size=10, color=ACCENT)

add_text_box(slide, Inches(0.5), Inches(6.6), Inches(12), Inches(0.6),
             "Each template auto-populates risks, standards, budget categories, milestones, and timelines — "
             "teams are productive within minutes of project creation, not weeks.",
             font_size=14, color=DGRAY)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 21 — ROLE-BASED ACCESS & TIERS
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), W, Inches(1.1), NAVY)
add_text_box(slide, Inches(0.6), Inches(0.3), Inches(10), Inches(0.6),
             "Role-Based Access & Subscription Tiers",
             font_size=28, color=WHITE, bold=True)

# Roles table
roles = [
    ("PMP (Project Manager)", "Full access to all 16 tabs; FDA Comms exclusive", TEAL),
    ("Technology", "Technical/Regulatory tabs, Risks, Docs, Actions, Resources, Message Board", ACCENT),
    ("Business", "Business milestones, Budget, Investment, Cap Table, Message Board", GREEN),
    ("Accounting", "Budget, Cash/Runway, Cap Table (limited editing)", RGBColor(0x7B, 0x1F, 0xA2)),
]

for i, (role, access, color) in enumerate(roles):
    top = Inches(1.4) + i * Inches(0.65)
    add_shape_bg(slide, Inches(0.5), top, Inches(2.5), Inches(0.55), color)
    add_text_box(slide, Inches(0.6), top + Inches(0.08), Inches(2.3), Inches(0.4),
                 role, font_size=13, color=WHITE, bold=True)
    add_text_box(slide, Inches(3.2), top + Inches(0.08), Inches(5), Inches(0.4),
                 access, font_size=13, color=DGRAY)

# Tiers table
tiers = [
    ("Starter — $500/mo", "2 seats", "Dual-Track, Gates, Timeline, Budget", TEAL),
    ("Growth — $1,000/mo", "5 seats", "All except Cap Table, FDA Comms, US Investment", ACCENT),
    ("Scale — $2,000/mo", "10 seats", "All 16 tabs + FDA Comms + Cap Table", NAVY),
]

add_text_box(slide, Inches(0.5), Inches(4.1), Inches(5), Inches(0.5),
             "Subscription Tiers", font_size=22, color=NAVY, bold=True)

for i, (tier, seats, tabs, color) in enumerate(tiers):
    top = Inches(4.7) + i * Inches(0.7)
    add_shape_bg(slide, Inches(0.5), top, Inches(3.5), Inches(0.6), color)
    add_text_box(slide, Inches(0.6), top + Inches(0.1), Inches(3.3), Inches(0.4),
                 tier, font_size=14, color=WHITE, bold=True)
    add_text_box(slide, Inches(4.2), top + Inches(0.1), Inches(1.5), Inches(0.4),
                 seats, font_size=13, color=DGRAY, bold=True)
    add_text_box(slide, Inches(5.7), top + Inches(0.1), Inches(7), Inches(0.4),
                 tabs, font_size=13, color=DGRAY)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 22 — COMPETITIVE COMPARISON MATRIX
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), W, Inches(1.1), NAVY)
add_text_box(slide, Inches(0.6), Inches(0.3), Inches(10), Inches(0.6),
             "Competitive Comparison: Control Tower vs. Alternatives",
             font_size=28, color=WHITE, bold=True)

# Table headers
headers = ["Feature", "Control Tower", "Greenlight\nGuru", "Master\nControl",
           "Jama\nConnect", "Qualio", "Monday.com\n/ Asana"]
col_widths = [Inches(2.5), Inches(1.8), Inches(1.5), Inches(1.5),
              Inches(1.5), Inches(1.3), Inches(1.8)]
x_offset = Inches(0.5)

for j, (header, cw) in enumerate(zip(headers, col_widths)):
    left = x_offset
    add_shape_bg(slide, left, Inches(1.3), cw, Inches(0.6),
                 NAVY if j == 0 else (TEAL if j == 1 else LGRAY))
    add_text_box(slide, left + Inches(0.05), Inches(1.33), cw - Inches(0.1), Inches(0.55),
                 header, font_size=10,
                 color=WHITE if j <= 1 else DGRAY,
                 bold=True, alignment=PP_ALIGN.CENTER)
    x_offset += cw

features = [
    ("Dual-Track Milestones",      ["✅","❌","❌","❌","❌","❌"]),
    ("Phase-Gate System",           ["✅","❌","❌","❌","❌","❌"]),
    ("FDA Q-Sub Generator",        ["✅","❌","❌","❌","❌","❌"]),
    ("RTA Pre-Check",              ["✅","❌","❌","❌","❌","❌"]),
    ("ISO 14971 Risk Dashboard",   ["✅","✅","✅","❌","❌","❌"]),
    ("DHF/DMR Tracking",           ["✅","✅","✅","❌","✅","❌"]),
    ("CAPA Management",            ["✅","✅","✅","❌","✅","❌"]),
    ("Standards Compliance",       ["✅","⚠️","✅","✅","⚠️","❌"]),
    ("Budget & Cash Runway",       ["✅","❌","❌","❌","❌","❌"]),
    ("Investor Pipeline",          ["✅","❌","❌","❌","❌","❌"]),
    ("Cap Table",                  ["✅","❌","❌","❌","❌","❌"]),
    ("Trilingual (EN/CN/KO)",     ["✅","❌","❌","❌","❌","❌"]),
    ("21 CFR Part 11 Audit Trail", ["✅","✅","✅","⚠️","✅","❌"]),
    ("Role-Based Access",          ["✅","✅","✅","✅","✅","⚠️"]),
    ("Starting Price",             ["$500/mo","~$1K/mo","~$4K/mo","~$3K/mo","~$1K/mo","$10/seat"]),
]

for i, (feature, vals) in enumerate(features):
    top = Inches(1.9) + i * Inches(0.35)
    x_offset = Inches(0.5)
    row_bg = WHITE if i % 2 == 0 else LGRAY
    for j, cw in enumerate(col_widths):
        left = x_offset
        if j == 0:
            text = feature
            color = DGRAY
            fs = 9
            al = PP_ALIGN.LEFT
        else:
            text = vals[j - 1]
            color = GREEN if text == "✅" else (RED if text == "❌" else DGRAY)
            fs = 9
            al = PP_ALIGN.CENTER
        add_shape_bg(slide, left, top, cw, Inches(0.35), row_bg)
        add_text_box(slide, left + Inches(0.05), top + Inches(0.03),
                     cw - Inches(0.1), Inches(0.3), text,
                     font_size=fs, color=color, alignment=al)
        x_offset += cw


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 23 — CLOSING / CTA
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)

add_text_box(slide, Inches(1), Inches(1.2), Inches(11), Inches(1),
             "Why Control Tower?",
             font_size=42, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1.5), Inches(2.5), Inches(10), Inches(3), (
    "The only platform purpose-built for FDA 510(k) medical device development.\n\n"
    "✓  Replaces 6+ standalone tools (PM, QMS, risk, finance, investor, comms)\n"
    "✓  Automates FDA Q-Sub generation, RTA pre-checks, and MDUFA timeline tracking\n"
    "✓  Trilingual for US-China and US-Korea market corridors\n"
    "✓  21 CFR Part 11 audit trail from Day 1\n"
    "✓  7 device templates — productive in minutes, not weeks\n"
    "✓  Starting at $500/month — 10× cheaper than enterprise QMS platforms\n\n"
    "From Pre-Submission to FDA Clearance — one platform, one team, one source of truth."
), font_size=18, color=WHITE, alignment=PP_ALIGN.CENTER)

add_shape_bg(slide, Inches(4), Inches(5.9), Inches(5.3), Inches(0.03), TEAL)

add_text_box(slide, Inches(1), Inches(6.2), Inches(11), Inches(0.8),
             "510(k) Bridge  •  Gresham, Oregon  •  www.510kbridge.com\n"
             "Schedule a Demo  •  info@510kbridge.com",
             font_size=16, color=TEAL, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════════════════
prs.save(OUTPUT)
print(f"✅ Saved {prs.slides.__len__()} slides → {OUTPUT}")
