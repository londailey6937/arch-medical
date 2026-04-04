#!/usr/bin/env python3
"""Generate Control Tower Sales Pitch PowerPoint — CHINESE (中文) version."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Paths ────────────────────────────────────────────────────────────────
SCREENSHOTS = "/Users/londailey/510kBridge_Webpage/public/screenshots"
OUTPUT = "/Users/londailey/510kBridge_Webpage/Control_Tower_Sales_Deck_CN.pptx"

# ── Brand colours ────────────────────────────────────────────────────────
NAVY   = RGBColor(0x0B, 0x23, 0x3A)
TEAL   = RGBColor(0x00, 0x96, 0x88)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY  = RGBColor(0xF5, 0xF5, 0xF5)
DGRAY  = RGBColor(0x33, 0x33, 0x33)
ACCENT = RGBColor(0x00, 0x7A, 0xCC)
RED    = RGBColor(0xE5, 0x3E, 0x3E)
GREEN  = RGBColor(0x2E, 0x7D, 0x32)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height


# ── Helpers ──────────────────────────────────────────────────────────────
def add_bg(slide, color=NAVY):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_bg(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Microsoft YaHei"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_bullet_frame(slide, left, top, width, height, bullets, font_size=14,
                     color=DGRAY, title=None, title_size=18, title_color=NAVY):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    if title:
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(title_size)
        p.font.color.rgb = title_color
        p.font.bold = True
        p.font.name = "Microsoft YaHei"
        p.space_after = Pt(8)
        start = 1
    else:
        start = 0

    for i, bullet in enumerate(bullets):
        if i == 0 and not title:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Microsoft YaHei"
        p.space_after = Pt(4)
        p.level = 0
    return tf


def screenshot_path(name):
    # Prefer CN screenshot, fall back to EN
    cn_path = os.path.join(SCREENSHOTS, f"{name}-cn.png")
    en_path = os.path.join(SCREENSHOTS, f"{name}.png")
    if os.path.exists(cn_path):
        return cn_path
    if os.path.exists(en_path):
        return en_path
    return None


def add_screenshot(slide, name, left, top, width, height=None):
    path = screenshot_path(name)
    if not path:
        return None
    if height:
        pic = slide.shapes.add_picture(path, left, top, width, height)
    else:
        pic = slide.shapes.add_picture(path, left, top, width)
    return pic


def tab_slide(slide, tab_number, tab_name, group_name):
    """Common header band for every tab slide."""
    add_bg(slide, WHITE)
    # Top navy bar
    add_shape_bg(slide, Inches(0), Inches(0), W, Inches(1.1), NAVY)
    add_text_box(slide, Inches(0.6), Inches(0.15), Inches(8), Inches(0.5),
                 f"{group_name}",
                 font_size=13, color=TEAL, bold=True)
    add_text_box(slide, Inches(0.6), Inches(0.45), Inches(10), Inches(0.55),
                 f"标签页 {tab_number}: {tab_name}",
                 font_size=26, color=WHITE, bold=True)
    # Branding
    add_text_box(slide, Inches(10.5), Inches(0.3), Inches(2.5), Inches(0.5),
                 "510(k) Bridge", font_size=16, color=TEAL, bold=True,
                 alignment=PP_ALIGN.RIGHT)
    add_text_box(slide, Inches(10.5), Inches(0.6), Inches(2.5), Inches(0.4),
                 "Control Tower 控制塔", font_size=12, color=WHITE,
                 alignment=PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════════════════════════════════
# 幻灯片 1 — 封面
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)

add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1),
             "510(k) Bridge 控制塔",
             font_size=44, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(2.6), Inches(11), Inches(0.8),
             "首个专为FDA 510(k)医疗器械开发打造的项目管理平台",
             font_size=22, color=TEAL, bold=False, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(2), Inches(3.8), Inches(9), Inches(1.5),
             "18个集成标签页  •  三语支持 (EN / 中文 / 한국어)  •  基于角色的访问控制\n"
             "21 CFR Part 11审计追踪  •  ISO 14971风险管理  •  FDA Q-Sub自动化\n"
             "实时云同步  •  7个预配置设备模板  •  3个订阅层级",
             font_size=16, color=WHITE, alignment=PP_ALIGN.CENTER)

add_shape_bg(slide, Inches(4), Inches(5.6), Inches(5.3), Inches(0.03), TEAL)

add_text_box(slide, Inches(1), Inches(6.0), Inches(11), Inches(0.5),
             "510(k) Bridge  •  俄勒冈州格雷舍姆  •  www.510kbridge.com",
             font_size=14, color=WHITE, alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════
# 幻灯片 2 — 市场问题
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), W, Inches(1.1), NAVY)
add_text_box(slide, Inches(0.6), Inches(0.3), Inches(10), Inches(0.6),
             "痛点：市场上不存在FDA专用的项目管理平台",
             font_size=28, color=WHITE, bold=True)

add_bullet_frame(slide, Inches(0.8), Inches(1.4), Inches(5.8), Inches(5.5), [
    "▸ 医疗器械公司使用通用PM工具（Jira、Asana、Monday.com），对FDA法规完全不了解",
    "▸ 法规顾问依赖电子表格、邮件和经验知识",
    "▸ 没有现有平台能将技术里程碑与FDA法规门控同步",
    "▸ 510(k)提交由于文档混乱（DHF/DMR缺口）平均需要6-9个月的返工",
    "▸ 企业平均发送2-3次FDA Q-Sub才被接受——每次周期增加$10K-$25K延迟成本",
    "▸ ISO 14971风险管理在独立文档中跟踪，与项目决策脱节",
    "▸ 国际团队（中美、韩美通道）缺乏统一的多语言工作空间",
], font_size=15, color=DGRAY, title="为什么行业需要控制塔",
   title_size=22, title_color=NAVY)

add_bullet_frame(slide, Inches(7.0), Inches(1.4), Inches(5.8), Inches(5.5), [
    "▸ Greenlight Guru — 专注QMS，非项目管理；无里程碑/门控跟踪；价格昂贵",
    "▸ MasterControl — 企业级QMS；复杂昂贵（$50K+/年）；无510(k)工作流",
    "▸ Jama Connect — 仅需求管理；无预算、风险矩阵或FDA通讯功能",
    "▸ Arena Solutions (PTC) — 专注PLM/BOM；无法规提交工具",
    "▸ Qualio — 以文档为中心的QMS；无双轨制里程碑，无门控系统",
    "▸ Monday.com / Asana — 通用PM；无任何FDA术语或合规功能",
    "",
    "这些平台都不能提供从Pre-Sub到FDA许可的集成510(k)项目生命周期，"
    "也没有内置的FDA自动化、三语支持和基于角色的治理。",
], font_size=15, color=DGRAY, title="竞争格局差距",
   title_size=22, title_color=RED)

# ═══════════════════════════════════════════════════════════════════════
# 幻灯片 3 — 平台架构
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), W, Inches(1.1), NAVY)
add_text_box(slide, Inches(0.6), Inches(0.3), Inches(10), Inches(0.6),
             "平台架构：5大分组、18个集成标签页",
             font_size=28, color=WHITE, bold=True)

groups = [
    ("项目管理", TEAL, ["1. 双轨制里程碑", "2. 门控系统",
                        "3. 时间线", "4. 行动面板"]),
    ("法规监管", RGBColor(0xE6, 0x5C, 0x00), ["5. 法规跟踪器",
                  "6. 风险看板", "7. FDA通讯中心",
                  "17. 谓词查找器", "18. 指导文件搜索"]),
    ("文档管理", ACCENT, ["8. 文档控制", "9. 审计追踪"]),
    ("财务管理", RGBColor(0x7B, 0x1F, 0xA2), ["10. 预算", "11. 现金/跑道",
               "12. 美国投资", "13. 股权表"]),
    ("运营管理", GREEN, ["14. 资源", "15. 供应商",
                        "16. 消息板"]),
]

x_start = Inches(0.5)
y_start = Inches(1.5)
col_w = Inches(2.4)
gap = Inches(0.15)

for i, (name, color, tabs) in enumerate(groups):
    left = x_start + i * (col_w + gap)
    add_shape_bg(slide, left, y_start, col_w, Inches(0.5), color)
    add_text_box(slide, left, y_start + Inches(0.05), col_w, Inches(0.4),
                 name, font_size=13, color=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_shape_bg(slide, left, y_start + Inches(0.5), col_w, Inches(3.0), LGRAY)
    add_bullet_frame(slide, left + Inches(0.15), y_start + Inches(0.6),
                     col_w - Inches(0.3), Inches(2.8), tabs,
                     font_size=13, color=DGRAY)

add_shape_bg(slide, Inches(0.5), Inches(5.3), Inches(12.3), Inches(0.05), TEAL)
add_bullet_frame(slide, Inches(0.5), Inches(5.5), Inches(12.3), Inches(1.5), [
    "三语 (EN/中文/한국어)  •  基于角色的访问控制 (PMP / 技术 / 商务 / 财务)  •  "
    "3个订阅层级 ($500–$2,000/月)  •  7个设备模板  •  Supabase云+离线  •  "
    "21 CFR Part 11审计追踪  •  设置向导  •  实时协作"
], font_size=14, color=DGRAY, title="跨平台核心能力",
   title_size=18, title_color=NAVY)


# ═══════════════════════════════════════════════════════════════════════
# 标签页详情幻灯片 — 每个标签页一张，含截图+销售价值+竞品差距
# ═══════════════════════════════════════════════════════════════════════

tab_data = [
    (1, "双轨制里程碑", "项目管理", "dual-track", [
        "▸ 技术与法规里程碑并行跟踪——唯一能并排展示两个轨道的平台",
        "▸ 技术轨道：设计冻结→原型→台架验证→V&V→设计转移",
        "▸ 法规轨道：Pre-Sub Q会议→510(k)准备→提交→预期许可",
        "▸ 根据项目周期（通常18-23个月）自动生成时间线",
        "▸ 状态循环（未开始→进行中→完成/受阻）并完整记录到审计追踪",
        "▸ 每个里程碑可分配所有者（技术、法规或商务），确保责任明确",
    ],
     "Jira/Asana有任务板但没有技术+法规同步轨道的概念。"
     "Greenlight Guru跟踪QMS文档但不跟踪工程里程碑。"
     "没有竞品能将原型测试计划与FDA提交时间线关联。"
    ),

    (2, "门控系统", "项目管理", "gate-system", [
        "▸ 数字化阶段门控治理——每个门控有基于清单的审查标准",
        "▸ 根据项目长度自动生成2-6个门控",
        "▸ 门控决策：通过/不通过/有条件——PMP用时间戳记录",
        "▸ 技术和商务团队的利益相关者输入纳入门控决策",
        "▸ 门控笔记记录讨论要点、条件和承诺",
        "▸ 门控状态随标准完成自动更新（未开始→待审查→已批准→受阻）",
    ],
     "Stage-Gate®是众所周知的方法论，但没有SaaS工具为医疗器械实施它。"
     "通用PM工具需要手动电子表格门控。"
     "控制塔的门控系统直接引用510(k)交付物——"
     "例如，'DHF完成？'作为与DHF跟踪器链接的实际门控标准。"
    ),

    (3, "时间线", "项目管理", "timeline", [
        "▸ 按月整合的项目视图，显示技术+商务活动",
        "▸ 影响指标：中性(|)、警告(⚠)、严重(🔴)用于进度风险标识",
        "▸ 从设备模板或向导输入自动生成——即时项目骨架",
        "▸ 可编辑事件，团队可修改文本和影响级别",
        "▸ 事件与双轨制里程碑关联，实现完整可追溯性",
        "▸ 为管理层提供整个510(k)项目生命周期的单页视图",
    ],
     "MS Project或Monday.com中的甘特图是通用进度条——"
     "它们无法区分FDA里程碑和营销活动。"
     "控制塔的时间线专为510(k)项目打造，"
     "内置FDA术语和影响标记。"
    ),

    (4, "行动面板", "项目管理", "actions", [
        "▸ 完整任务板：标题、负责人、优先级(高/中/低)、状态(待办/进行中/完成/受阻)、截止日期",
        "▸ 门控关联操作——每个任务可引用特定门控(G1、G2等)",
        "▸ DHF文档跟踪器——21 CFR 820.30设计历史文件文档清单",
        "▸ DMR文档跟踪器——21 CFR 820.181设备主记录（12种文档类型）",
        "▸ CAPA日志——纠正与预防措施，与风险看板的风险关联",
        "▸ 一个标签页中四个集成子板块：任务+DHF+DMR+CAPA",
    ],
     "Greenlight Guru有DHF跟踪但与项目管理任务隔离。"
     "Jira有任务板但没有DHF/DMR/CAPA意识。"
     "控制塔独特地将任务管理与FDA文档清单结合，"
     "并将CAPA直接链接到风险看板中识别的风险。"
    ),

    (5, "法规跟踪器", "法规监管", "regulatory-tracker", [
        "▸ 标准合规仪表板，从设备模板自动填充",
        "▸ 跟踪IEC 60601-1、IEC 60601-1-2、ISO 10993、ISO 14971、IEC 62304、21 CFR 820、ISO 13485",
        "▸ 条款级跟踪——深入到各子条款(如ISO 10993-5、-10)",
        "▸ 每个标准0-100%进度条，含状态循环",
        "▸ 证据链接——为每个条款附加文档引用",
        "▸ 即时显示哪些标准落后，集中团队精力",
    ],
     "Qualio和MasterControl跟踪文档合规但不跟踪标准级进度。"
     "没有竞品能根据设备类型自动填充适用标准。"
     "团队目前在Excel中维护合规矩阵——"
     "控制塔消除了那个电子表格并使其实时化。"
    ),

    (6, "风险看板", "法规监管", "risk-dashboard", [
        "▸ 完整的ISO 14971风险管理可视化，含严重性×概率矩阵",
        "▸ 颜色编码风险等级：绿色(可接受)→黄色(ALARP)→红色(不可接受)",
        "▸ 每个风险字段：严重性、概率、控制措施、残余风险、缓解状态、模块、标准",
        "▸ 红色风险警报自动在FDA通讯标签页中触发警告",
        "▸ 模板预填充风险(每类设备5-8个——如呼吸设备含EMG特定风险)",
        "▸ 点击风险ID打开详细编辑器——所有更改记录到审计追踪",
    ],
     "大多数公司在Word文档或独立Excel矩阵中跟踪风险。"
     "Greenlight Guru有风险管理但与项目时间线脱钩。"
     "控制塔的风险看板将风险与CAPA、门控标准"
     "和FDA Pre-Sub问题连接——形成闭环风险系统。"
    ),

    (7, "FDA通讯中心", "法规监管", "fda-comms", [
        "▸ PMP独占标签页——FDA互动的神经中枢",
        "▸ Q-Sub封面信生成器——5种类型(Pre-Sub会议、书面问题、受控通信等)",
        "▸ 导出问题包——从消息板的Pre-Sub问题线程提取，格式化为FDA包",
        "▸ 拒绝接受(RTA)清单——17项自检对照FDA接受标准",
        "▸ MDUFA审查时间线——标准510(k)审查的第1天至第90天里程碑跟踪",
        "▸ 实质等效(SE)决策流程图——FDA的SE判定逻辑可视化",
    ],
     "没有竞品有此功能。没有现有平台能生成FDA Q-Sub封面信、"
     "导出Pre-Sub问题包或提供RTA预检。"
     "这是控制塔最强的差异化优势——"
     "它自动化了最耗时和最容易出错的法规沟通工作流。"
    ),

    (8, "文档控制", "文档管理", "document-control", [
        "▸ ISO 13485对齐的文档生命周期：草稿→审核中→已批准→生效→已废弃",
        "▸ 自动生成文档控制编号(DCN-REG-001等)",
        "▸ 类别：法规、技术、商务、法律、财务、模板",
        "▸ 浏览器本地存储保护知识产权+已批准文档可选同步到服务器",
        "▸ 源引用链接(GitHub commit、SVN revision等)",
        "▸ 里程碑关联——例如R8里程碑触发DHF文档就绪检查",
    ],
     "MasterControl和Qualio是完整QMS平台，费用$50K+/年。"
     "控制塔以极低成本提供80%的关键文档控制功能，"
     "与项目里程碑紧密集成，"
     "而不是作为独立系统存在。"
    ),

    (9, "审计追踪", "文档管理", "audit-trail", [
        "▸ 符合21 CFR Part 11的不可变活动日志",
        "▸ 每次更改均记录：时间戳、用户角色、操作类型、目标、旧值→新值",
        "▸ 100+种操作类型(里程碑状态、风险字段、门控决策、文档状态、预算条目等)",
        "▸ 关键词搜索和操作类型过滤",
        "▸ 导出完整历史记录为CSV/JSON，用于法规提交",
        "▸ 自动同步到Supabase后端；离线时本地排队——零数据丢失",
    ],
     "这是受监管行业的基本要求，但通用PM工具缺少此功能。"
     "Jira的审计日志面向IT而非法规。"
     "控制塔的审计追踪专为FDA调查人员和认证机构审核员设计，"
     "捕获他们期望的确切字段和格式。"
    ),

    (10, "预算", "财务管理", "budget", [
        "▸ 按类别的预算对比实际跟踪——从设备模板预填充",
        "▸ 类别：原型与材料、测试、法规与法律、人员等",
        "▸ 每类字段：计划金额、实际金额、差异(自动计算)、备注",
        "▸ 双币种支持：美元/人民币自动切换(7.25×汇率)",
        "▸ 随项目演进添加/编辑/删除类别",
        "▸ 所有预算变更记录到审计追踪，确保财务问责",
    ],
     "通用PM工具要么完全没有预算功能（Jira），要么提供基本的支出跟踪"
     "但没有医疗器械特定的成本类别。"
     "控制塔预填充与510(k)开发阶段匹配的预算类别——"
     "测试实验室费用、法规顾问费、生物相容性研究费等。"
    ),

    (11, "现金/跑道", "财务管理", "cash-runway", [
        "▸ 实时财务健康仪表板：现金余额、月度消耗率、跑道(月数)",
        "▸ 跑道颜色编码：绿色(>12个月)、黄色(6-12个月)、红色(<6个月)——即时可见",
        "▸ 融资轮次跟踪：管线→已承诺→已收到，含金额和日期",
        "▸ 按月消耗历史：实际对比计划支出",
        "▸ 可视化消耗历史图表，用于董事会演示",
        "▸ 对初创企业至关重要——回答'我们的资金能撑到FDA许可吗？'",
    ],
     "财务规划工具（QuickBooks、Xero）不将现金跑道与法规里程碑关联。"
     "初创仪表板（Carta、Visible）跟踪消耗但对FDA时间线一无所知。"
     "控制塔独特地回答了每个医疗器械创始人都会问的问题："
     "'我的钱能撑到510(k)许可吗？'"
    ),

    (12, "美国投资", "财务管理", "us-investment", [
        "▸ 投资者关系管道：潜在→已联系→洽谈中→条款书→已承诺",
        "▸ 投资者档案：类型(VC、天使、战略、PE、政府)、阶段匹配、承诺金额",
        "▸ IR活动记录：会议、演示、尽职调查、条款谈判",
        "▸ 活动指标和转化率，评估筹资效果",
        "▸ 专为进入美国市场的国际医疗器械公司打造",
        "▸ 整合原本散落在电子表格和CRM中的投资者跟踪",
    ],
     "CRM工具（HubSpot、Salesforce）跟踪投资者联系但没有医疗器械背景。"
     "Carta处理股权表但不处理投资者管道。"
     "控制塔提供专业的投资者管道，与510(k)项目并排——"
     "创始人可以向投资者准确展示他们在法规流程中的位置。"
    ),

    (13, "股权表", "财务管理", None, [
        "▸ 完整的股权所有权登记：普通股、优先股A/B/C、期权、认股权证",
        "▸ 自动计算所有权百分比和董事会席位跟踪",
        "▸ 股权事件：融资轮次、股票分割、期权授予、优先股转换、认股权证行使",
        "▸ 归属计划管理：4年标准+1年悬崖期(可自定义)",
        "▸ 当前已归属/未归属跟踪，含下次归属日期计算",
        "▸ 三个集成部分：股东+股权事件+归属计划",
    ],
     "Carta是股权表的黄金标准，但费用$5K-$50K/年且是独立工具。"
     "控制塔的股权表与项目集成——"
     "当融资轮次关闭（美国投资标签页），可以流入股权表。"
     "对于早期医疗器械初创企业，又消除了一个独立工具。"
    ),

    (14, "资源", "运营管理", "resources", [
        "▸ 跨工作流的团队分配和利用率监控",
        "▸ 工作流分配：技术、法规、商务、财务(+自定义)的百分比",
        "▸ 利用率仪表：绿色(<85%)、黄色(85-100%)、红色(>100%)——防止过劳",
        "▸ 团队成员档案：姓名、角色、邮箱、简介",
        "▸ 内联编辑——点击百分比实时调整分配",
        "▸ 所有分配变更记录到审计追踪，确保人力规划透明",
    ],
     "Jira或Monday.com的资源管理是基于任务小时的，而非工作流百分比。"
     "医疗器械团队经常有身兼多职的工程师（技术+法规）。"
     "控制塔的分配模型匹配小型设备团队的实际工作方式——"
     "一目了然地显示V&V工程师是否在法规文档方面也过载。"
    ),

    (15, "供应商", "运营管理", "suppliers", [
        "▸ 供应商资质审核和PO跟踪——符合21 CFR 820供应商控制要求",
        "▸ 状态生命周期：审查中→已合格→活跃→暂停→已拒绝",
        "▸ 字段：组件/零件号、交期(天)、PO状态、合同制造里程碑",
        "▸ 资质标准：ISO 13485/9001认证、测试报告、业务连续性计划",
        "▸ 对510(k)至关重要——FDA期望在QMS中有文档化的供应商资质",
        "▸ 与项目里程碑关联——当组件交期威胁门控评审时发出提醒",
    ],
     "Arena Solutions (PTC)是完整PLM加BOM管理，但费用$1,000+/席位/年。"
     "控制塔提供FDA期望的供应商资质跟踪，"
     "无需完整PLM的复杂性——完美适合"
     "首次用5-15个供应商构建设备的初创企业。"
    ),

    (16, "消息板", "通讯协作", "message-board", [
        "▸ 跨职能线程讨论，含决策和行动跟踪",
        "▸ 工作流：技术、法规、商务、财务、Pre-Sub问题、其他",
        "▸ 意图标签：讨论、决策、通知、升级——每个对话都有目的",
        "▸ [决策]和[行动]消息前缀，确保可追溯性",
        "▸ Pre-Sub问题工作流：协作编写FDA问题→导出到FDA通讯中心",
        "▸ 视图：全部线程、我的项目、决策、管理层——按生命周期过滤(打开/已解决)",
    ],
     "Slack、Teams和电子邮件是医疗器械决策消亡的地方——"
     "关键法规决策埋没在没有可追溯性的聊天流中。"
     "控制塔的消息板捕获每个决策的角色归属、时间戳，"
     "以及生成行动项的能力——"
     "然后将Pre-Sub问题直接导出到FDA通讯中心。"
    ),

    (17, "谓词查找器", "法规监管", None, [
        "▸ 基于openFDA的独立SaaS工具——搜索200K+已许可的510(k)设备",
        "▸ 按产品代码、申请人、设备名称或K编号搜索",
        "▸ 实时FDA数据库查询——始终是最新数据，无过时导出",
        "▸ 一键选择谓词设备——直接导入您的510(k)项目",
        "▸ 每个结果显示决策编号、许可日期、产品代码和声明链接",
        "▸ 免费层级访问——即使没有控制塔订阅也可使用",
    ],
     "法规团队目前通过笨重的FDA.gov搜索手动查找510(k)数据库。"
     "没有竞品将谓词设备搜索直接集成到510(k)工作流中。"
     "控制塔的谓词查找器自动化了510(k)最关键的第一步——"
     "识别您的谓词设备——并将其直接链接到项目中。"
    ),

    (18, "FDA指导文件搜索", "法规监管", None, [
        "▸ 与510(k)提交相关的FDA指导文件可搜索数据库",
        "▸ 按设备类型、主题领域、发布日期和草案/最终状态过滤",
        "▸ 快速参考摘要——无需阅读100页PDF即可了解指导适用性",
        "▸ 链接到官方FDA.gov源文件，便于在提交中引用",
        "▸ 指导收藏夹——保存常用参考文件以便快速访问",
        "▸ 根据设备模板和适用标准自动推荐",
    ],
     "法规专业人员花费数小时在FDA.gov上搜索相关指导文件。"
     "没有竞品平台在项目上下文中索引或呈现FDA指导。"
     "控制塔将指导发现置于需要的工作流中——"
     "消除了减慢法规策略的标签页切换和手动搜索。"
    ),
]

for tab_num, tab_name, group, screenshot_key, value_props, competitor_gap in tab_data:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tab_slide(slide, tab_num, tab_name, group)

    ss = None
    if screenshot_key:
        ss = add_screenshot(slide, screenshot_key,
                            Inches(0.4), Inches(1.3), Inches(6.2))

    if ss:
        text_left = Inches(6.9)
        text_width = Inches(6.0)
    else:
        text_left = Inches(0.8)
        text_width = Inches(11.5)

    add_bullet_frame(slide, text_left, Inches(1.3), text_width, Inches(3.5),
                     value_props, font_size=12, color=DGRAY,
                     title="价值主张", title_size=16, title_color=TEAL)

    add_shape_bg(slide, Inches(0.3), Inches(5.6), Inches(12.7), Inches(1.6), LGRAY)
    add_text_box(slide, Inches(0.5), Inches(5.65), Inches(3), Inches(0.35),
                 "竞品差距分析", font_size=13, color=RED, bold=True)
    add_text_box(slide, Inches(0.5), Inches(5.95), Inches(12.3), Inches(1.2),
                 competitor_gap, font_size=12, color=DGRAY)


# ═══════════════════════════════════════════════════════════════════════
# 跨平台功能幻灯片
# ═══════════════════════════════════════════════════════════════════════

# ── 利益相关者输入与变更请求 ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), W, Inches(1.1), NAVY)
add_text_box(slide, Inches(0.6), Inches(0.15), Inches(8), Inches(0.5),
             "跨平台功能",
             font_size=13, color=TEAL, bold=True)
add_text_box(slide, Inches(0.6), Inches(0.45), Inches(10), Inches(0.55),
             "利益相关者输入与变更请求",
             font_size=26, color=WHITE, bold=True)
add_text_box(slide, Inches(10.5), Inches(0.3), Inches(2.5), Inches(0.5),
             "510(k) Bridge", font_size=16, color=TEAL, bold=True,
             alignment=PP_ALIGN.RIGHT)

add_bullet_frame(slide, Inches(0.8), Inches(1.3), Inches(5.8), Inches(3.5), [
    "▸ 每个门控系统页面上的浮动📥按钮——打开侧面板",
    "▸ 技术和商务利益相关者按门控提交结构化输入",
    "▸ 字段：门控选择、输入类型(技术/商务/风险/资源/其他)、自由文本备注",
    "▸ PMP在门控详情视图中查看所有提交——按贡献者组织",
    "▸ 输入直接纳入门控决策——无需翻找邮件/Slack",
    "▸ 确保在门控评审前捕获跨职能意见",
], font_size=13, color=DGRAY, title="利益相关者输入",
   title_size=18, title_color=TEAL)

add_bullet_frame(slide, Inches(7.0), Inches(1.3), Inches(5.8), Inches(3.5), [
    "▸ 非PMP团队成员可通过结构化变更请求表请求变更",
    "▸ 涵盖：里程碑日期、预算调整、风险重新评估、资源重新分配",
    "▸ PMP审查并批准/拒绝，带时间戳、理由和审计追踪",
    "▸ 防止未经授权的变更，同时保持团队敏捷性",
    "▸ 每个请求和决策都记录在21 CFR Part 11审计追踪中",
    "▸ 平衡治理与速度——对受监管产品开发至关重要",
], font_size=13, color=DGRAY, title="变更请求工作流",
   title_size=18, title_color=TEAL)

add_shape_bg(slide, Inches(0.3), Inches(5.3), Inches(12.7), Inches(1.8), LGRAY)
add_bullet_frame(slide, Inches(0.5), Inches(5.4), Inches(12.3), Inches(1.6), [
    "没有竞品提供绑定到阶段门控决策的结构化利益相关者输入。在通用PM工具中，"
    "利益相关者反馈分散在评论、邮件和会议中，没有可追溯性。"
    "控制塔确保每个声音都被听到，每个变更都受治理，每个决策都有审计日志。",
], font_size=13, color=DGRAY, title="为什么重要",
   title_size=15, title_color=RED)

# ── QMS-Lite初创版 ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), W, Inches(1.1), NAVY)
add_text_box(slide, Inches(0.6), Inches(0.15), Inches(8), Inches(0.5),
             "跨平台功能",
             font_size=13, color=TEAL, bold=True)
add_text_box(slide, Inches(0.6), Inches(0.45), Inches(10), Inches(0.55),
             "QMS-Lite 初创企业质量管理",
             font_size=26, color=WHITE, bold=True)
add_text_box(slide, Inches(10.5), Inches(0.3), Inches(2.5), Inches(0.5),
             "510(k) Bridge", font_size=16, color=TEAL, bold=True,
             alignment=PP_ALIGN.RIGHT)

qms_mappings = [
    ("文档控制 (ISO 13485 §4.2.4)", "→ 标签页8：文档控制"),
    ("风险管理 (ISO 14971)", "→ 标签页6：风险看板"),
    ("设计与开发控制 (21 CFR 820.30)", "→ 标签页1-4：里程碑+行动项+DHF/DMR"),
    ("CAPA (21 CFR 820.198)", "→ 标签页4：行动面板——CAPA日志"),
    ("供应商控制 (21 CFR 820.50)", "→ 标签页15：供应商"),
    ("审计追踪 (21 CFR Part 11)", "→ 标签页9：审计追踪"),
    ("管理评审", "→ 标签页2：门控系统——门控决策"),
    ("培训记录", "→ 标签页14：资源——团队资质"),
    ("法规提交", "→ 标签页7：FDA通讯中心"),
]

for i, (req, mapping) in enumerate(qms_mappings):
    top = Inches(1.4) + i * Inches(0.48)
    add_shape_bg(slide, Inches(0.5), top, Inches(5.5), Inches(0.42),
                 LGRAY if i % 2 == 0 else WHITE)
    add_text_box(slide, Inches(0.6), top + Inches(0.05), Inches(5.3), Inches(0.35),
                 req, font_size=12, color=DGRAY, bold=True)
    add_shape_bg(slide, Inches(6.0), top, Inches(6.5), Inches(0.42),
                 LGRAY if i % 2 == 0 else WHITE)
    add_text_box(slide, Inches(6.1), top + Inches(0.05), Inches(6.3), Inches(0.35),
                 mapping, font_size=12, color=TEAL)

add_shape_bg(slide, Inches(0.3), Inches(5.8), Inches(12.7), Inches(1.4), LGRAY)
add_bullet_frame(slide, Inches(0.5), Inches(5.9), Inches(12.3), Inches(1.2), [
    "初创企业不需要Greenlight Guru ($30K+/年)或MasterControl ($50K+/年)来证明QMS就绪性。"
    "控制塔将9个核心QMS要求映射到现有标签页——以极低成本为初创企业提供80% QMS解决方案。"
    "非常适合准备首次FDA提交的种子轮公司。",
], font_size=13, color=DGRAY, title="以初创企业友好的价格获得80% QMS覆盖",
   title_size=15, title_color=NAVY)

# ── 通知、警报与导出报告 ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), W, Inches(1.1), NAVY)
add_text_box(slide, Inches(0.6), Inches(0.15), Inches(8), Inches(0.5),
             "跨平台功能",
             font_size=13, color=TEAL, bold=True)
add_text_box(slide, Inches(0.6), Inches(0.45), Inches(10), Inches(0.55),
             "通知、警报与高管报告",
             font_size=26, color=WHITE, bold=True)
add_text_box(slide, Inches(10.5), Inches(0.3), Inches(2.5), Inches(0.5),
             "510(k) Bridge", font_size=16, color=TEAL, bold=True,
             alignment=PP_ALIGN.RIGHT)

add_bullet_frame(slide, Inches(0.8), Inches(1.3), Inches(5.8), Inches(3.5), [
    "▸ 通知栏含5种警报类型：逾期里程碑、门控评审到期、"
    "高优先级行动项、预算偏差和红色风险警报",
    "▸ 通知铃铛上的徽章计数——一目了然查看未读警报",
    "▸ 点击导航：每个警报直接链接到相关标签页和记录",
    "▸ 按用户角色可配置的通知偏好",
    "▸ 主动警报防止错过截止日期和门控延迟",
], font_size=13, color=DGRAY, title="智能通知与警报",
   title_size=18, title_color=TEAL)

add_bullet_frame(slide, Inches(7.0), Inches(1.3), Inches(5.8), Inches(3.5), [
    "▸ 一键高管报告——整个项目状态的综合HTML导出",
    "▸ 涵盖：里程碑进度、门控状态、风险摘要、预算偏差、团队利用率",
    "▸ 专业格式，可直接用于董事会演示或投资者更新",
    "▸ 包含法规就绪性护栏：3面板系统含11项自动检查",
    "▸ 检查项：SE论证强度、性能数据完整性、生物相容性状态、"
    "软件文档(IEC 62304)、21 CFR 820设计控制等",
    "▸ 自动生成护栏评分，为团队提供提交前信心评级",
], font_size=13, color=DGRAY, title="导出报告与就绪性护栏",
   title_size=18, title_color=TEAL)

add_shape_bg(slide, Inches(0.3), Inches(5.3), Inches(12.7), Inches(1.8), LGRAY)
add_bullet_frame(slide, Inches(0.5), Inches(5.4), Inches(12.3), Inches(1.6), [
    "通用PM工具有基本通知但零法规感知。"
    "控制塔的警报具备FDA上下文感知——它知道门控评审何时逾期、"
    "红色风险何时缺少缓解措施、Pre-Sub截止日期何时临近。"
    "导出报告取代了每月消耗4-8小时的手动PowerPoint状态报告。",
], font_size=13, color=DGRAY, title="为什么重要",
   title_size=15, title_color=RED)


# ═══════════════════════════════════════════════════════════════════════
# 幻灯片 — 设备模板与向导
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), W, Inches(1.1), NAVY)
add_text_box(slide, Inches(0.6), Inches(0.3), Inches(10), Inches(0.6),
             "设置向导与7个预配置设备模板",
             font_size=28, color=WHITE, bold=True)

templates = [
    ("🫁 呼吸设备", "呼吸机、CPAP、雾化器\n产品代码：BZD, CBK, BTT, FRA"),
    ("❤️ 心血管设备", "心电监护仪、血压设备、支架\n产品代码：DRX, DXH, MHX"),
    ("🦴 骨科设备", "关节植入物、器械、\n固定装置"),
    ("🧪 IVD体外诊断", "体外诊断、检测试剂、\n分析仪"),
    ("📷 影像设备", "X射线、超声波、\nMRI附件"),
    ("♿ 康复设备", "治疗设备、\n移动辅助器"),
    ("💻 SaMD软件", "作为医疗器械的软件\n(IEC 62304生命周期)"),
]

for i, (name, desc) in enumerate(templates):
    col = i % 4
    row = i // 4
    left = Inches(0.5) + col * Inches(3.15)
    top = Inches(1.5) + row * Inches(2.5)
    add_shape_bg(slide, left, top, Inches(2.9), Inches(2.2), LGRAY)
    add_text_box(slide, left + Inches(0.15), top + Inches(0.1),
                 Inches(2.6), Inches(0.4), name,
                 font_size=16, color=NAVY, bold=True)
    add_text_box(slide, left + Inches(0.15), top + Inches(0.55),
                 Inches(2.6), Inches(0.7), desc,
                 font_size=11, color=DGRAY)
    add_text_box(slide, left + Inches(0.15), top + Inches(1.3),
                 Inches(2.6), Inches(0.8),
                 "包含：\n• 预填充的风险与标准\n• 预算类别\n• 时间线与里程碑",
                 font_size=10, color=ACCENT)

add_text_box(slide, Inches(0.5), Inches(6.6), Inches(12), Inches(0.6),
             "每个模板自动填充风险、标准、预算类别、里程碑和时间线——"
             "团队在创建项目后几分钟内即可高效工作，而非数周。",
             font_size=14, color=DGRAY)


# ═══════════════════════════════════════════════════════════════════════
# 幻灯片 21 — 角色与层级
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), W, Inches(1.1), NAVY)
add_text_box(slide, Inches(0.6), Inches(0.3), Inches(10), Inches(0.6),
             "基于角色的访问与订阅层级",
             font_size=28, color=WHITE, bold=True)

roles = [
    ("PMP (项目经理)", "完全访问所有18个标签页；FDA通讯独占权限", TEAL),
    ("技术", "技术/法规标签页、风险、文档、行动项、资源、消息板", ACCENT),
    ("商务", "商务里程碑、预算、投资、股权表、消息板", GREEN),
    ("财务", "预算、现金/跑道、股权表（有限编辑）", RGBColor(0x7B, 0x1F, 0xA2)),
]

for i, (role, access, color) in enumerate(roles):
    top = Inches(1.4) + i * Inches(0.65)
    add_shape_bg(slide, Inches(0.5), top, Inches(2.5), Inches(0.55), color)
    add_text_box(slide, Inches(0.6), top + Inches(0.08), Inches(2.3), Inches(0.4),
                 role, font_size=13, color=WHITE, bold=True)
    add_text_box(slide, Inches(3.2), top + Inches(0.08), Inches(5), Inches(0.4),
                 access, font_size=13, color=DGRAY)

tiers = [
    ("入门版 — $500/月", "2个席位", "双轨制、门控、时间线、预算", TEAL),
    ("成长版 — $1,000/月", "5个席位", "除股权表、FDA通讯、美国投资外全部", ACCENT),
    ("规模版 — $2,000/月", "10个席位", "全部18个标签页 + FDA通讯 + 股权表 + 谓词查找器", NAVY),
]

add_text_box(slide, Inches(0.5), Inches(4.1), Inches(5), Inches(0.5),
             "订阅层级", font_size=22, color=NAVY, bold=True)

for i, (tier, seats, tabs, color) in enumerate(tiers):
    top = Inches(4.7) + i * Inches(0.7)
    add_shape_bg(slide, Inches(0.5), top, Inches(3.5), Inches(0.6), color)
    add_text_box(slide, Inches(0.6), top + Inches(0.1), Inches(3.3), Inches(0.4),
                 tier, font_size=14, color=WHITE, bold=True)
    add_text_box(slide, Inches(4.2), top + Inches(0.1), Inches(1.5), Inches(0.4),
                 seats, font_size=13, color=DGRAY, bold=True)
    add_text_box(slide, Inches(5.7), top + Inches(0.1), Inches(7), Inches(0.4),
                 tabs, font_size=13, color=DGRAY)


# ═══════════════════════════════════════════════════════════════════════
# 幻灯片 22 — 竞品对比矩阵
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), W, Inches(1.1), NAVY)
add_text_box(slide, Inches(0.6), Inches(0.3), Inches(10), Inches(0.6),
             "竞品对比：控制塔 vs. 替代方案",
             font_size=28, color=WHITE, bold=True)

headers = ["功能", "控制塔", "Greenlight\nGuru", "Master\nControl",
           "Jama\nConnect", "Qualio", "Monday.com\n/ Asana"]
col_widths = [Inches(2.5), Inches(1.8), Inches(1.5), Inches(1.5),
              Inches(1.5), Inches(1.3), Inches(1.8)]
x_offset = Inches(0.5)

for j, (header, cw) in enumerate(zip(headers, col_widths)):
    left = x_offset
    add_shape_bg(slide, left, Inches(1.3), cw, Inches(0.6),
                 NAVY if j == 0 else (TEAL if j == 1 else LGRAY))
    add_text_box(slide, left + Inches(0.05), Inches(1.33), cw - Inches(0.1), Inches(0.55),
                 header, font_size=10,
                 color=WHITE if j <= 1 else DGRAY,
                 bold=True, alignment=PP_ALIGN.CENTER)
    x_offset += cw

features = [
    ("双轨制里程碑",               ["✅","❌","❌","❌","❌","❌"]),
    ("阶段门控系统",               ["✅","❌","❌","❌","❌","❌"]),
    ("FDA Q-Sub生成器",            ["✅","❌","❌","❌","❌","❌"]),
    ("RTA预检查",                  ["✅","❌","❌","❌","❌","❌"]),
    ("谓词查找器 (openFDA)",       ["✅","❌","❌","❌","❌","❌"]),
    ("FDA指导文件搜索",            ["✅","❌","❌","❌","❌","❌"]),
    ("ISO 14971风险看板",          ["✅","✅","✅","❌","❌","❌"]),
    ("DHF/DMR跟踪",               ["✅","✅","✅","❌","✅","❌"]),
    ("CAPA管理",                   ["✅","✅","✅","❌","✅","❌"]),
    ("利益相关者门控输入",         ["✅","❌","❌","❌","❌","❌"]),
    ("预算与现金跑道",             ["✅","❌","❌","❌","❌","❌"]),
    ("投资者管道 + 股权表",        ["✅","❌","❌","❌","❌","❌"]),
    ("QMS-Lite (9项要求)",         ["✅","✅","✅","❌","✅","❌"]),
    ("三语 (EN/中文/한국어)",      ["✅","❌","❌","❌","❌","❌"]),
    ("21 CFR Part 11审计追踪",     ["✅","✅","✅","⚠️","✅","❌"]),
    ("起价",                       ["$500/月","~$1K/月","~$4K/月","~$3K/月","~$1K/月","$10/席"]),
]

for i, (feature, vals) in enumerate(features):
    top = Inches(1.9) + i * Inches(0.33)
    x_offset = Inches(0.5)
    row_bg = WHITE if i % 2 == 0 else LGRAY
    for j, cw in enumerate(col_widths):
        left = x_offset
        if j == 0:
            text = feature
            color = DGRAY
            fs = 9
            al = PP_ALIGN.LEFT
        else:
            text = vals[j - 1]
            color = GREEN if text == "✅" else (RED if text == "❌" else DGRAY)
            fs = 9
            al = PP_ALIGN.CENTER
        add_shape_bg(slide, left, top, cw, Inches(0.33), row_bg)
        add_text_box(slide, left + Inches(0.05), top + Inches(0.02),
                     cw - Inches(0.1), Inches(0.29), text,
                     font_size=fs, color=color, alignment=al)
        x_offset += cw


# ═══════════════════════════════════════════════════════════════════════
# 幻灯片 23 — 结语/行动号召
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)

add_text_box(slide, Inches(1), Inches(0.8), Inches(11), Inches(1),
             "为什么选择控制塔？",
             font_size=42, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1.5), Inches(1.9), Inches(10), Inches(4.2), (
    "唯一专为FDA 510(k)医疗器械开发打造的平台。\n\n"
    "✓  替代6+个独立工具（PM、QMS、风险、财务、投资、通讯）\n"
    "✓  18个集成标签页，包括谓词查找器和FDA指导文件搜索\n"
    "✓  自动化FDA Q-Sub生成、RTA预检查和MDUFA时间线跟踪\n"
    "✓  利益相关者输入、变更请求与法规就绪性护栏\n"
    "✓  QMS-Lite映射9项核心质量要求——以初创企业价格获得80% QMS\n"
    "✓  三语支持，服务中美和韩美市场通道\n"
    "✓  从第一天起具备21 CFR Part 11审计追踪  •  7个设备模板\n"
    "✓  起价$500/月——比企业QMS平台便宜10倍\n\n"
    "从Pre-Submission到FDA许可——一个平台，一个团队，一个真实来源。"
), font_size=15, color=WHITE, alignment=PP_ALIGN.CENTER)

add_shape_bg(slide, Inches(4), Inches(6.2), Inches(5.3), Inches(0.03), TEAL)

add_text_box(slide, Inches(1), Inches(6.45), Inches(11), Inches(0.8),
             "510(k) Bridge  •  俄勒冈州格雷舍姆  •  www.510kbridge.com\n"
             "预约演示  •  info@510kbridge.com",
             font_size=16, color=TEAL, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════
# 保存
# ═══════════════════════════════════════════════════════════════════════
prs.save(OUTPUT)
print(f"✅ 已保存 {prs.slides.__len__()} 张幻灯片 → {OUTPUT}")
